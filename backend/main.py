from pptx import Presentation
import lxml.etree as etree
import uuid
import json
from typing import List
from core import extract_slide_mapping,prepare_sections,replace_slides, compile_sections

# folder="output"
# pres = Presentation('master.pptx')
# slidelist=pres.slides._sldIdLst
# mapping=extract_slide_mapping(slidelist)

# keys=['APM','SSO']
# l=prepare_sections(keys,pres,mapping)
# replace_slides(l,pres,folder,save=True)

def create_pptx(presentation: Presentation,sections: List[str]):
    folder="output"
    slidelist=presentation.slides._sldIdLst
    mapping=extract_slide_mapping(slidelist)


    new_xml=prepare_sections(sections,presentation,mapping)
    
    file_path=replace_slides(new_xml,presentation,folder,save=True)

    return file_path

def get_sections(presentation: Presentation):
    
    slidelist=presentation.slides._sldIdLst

    mapping = extract_slide_mapping(slidelist)

    res = compile_sections(presentation,mapping)
    res=list(res)

    section_dict=dict()
    section_content=list()

    for section in res:
        section_content.append(section)

    section_dict['count']=len(section_content)
    section_dict['data']=section_content


    return section_dict


if __name__ == "__main__":
    file_path="master.pptx"
    pres = Presentation('backend/master.pptx')
    
    sec = get_sections(pres)
    print(sec)